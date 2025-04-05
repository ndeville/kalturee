"""UPLOAD SINGLE POWERPOINT FILE TO KALTURA"""

from datetime import datetime
import os
import subprocess
import ollama
from pptx import Presentation

ts_time = f"{datetime.now().strftime('%H:%M:%S')}"
ts_file = f"{datetime.now().strftime('%y%m%d-%H%M')}"
print(f"\n---------- {ts_time} starting {os.path.basename(__file__)}")

# Start Chrono
import time
start_time = time.time()

# Credentials
from dotenv import load_dotenv
load_dotenv()
MY_USER_SECRET = os.getenv("user_secret")
MY_ADMIN_SECRET = os.getenv("admin_secret")
MY_PARTNER_ID = os.getenv("partner_id")

# Kaltura Client
from KalturaClient import *
from KalturaClient.Plugins.Core import *
# Import Document plugin
from KalturaClient.Plugins.Document import *




# THUMBNAIL (convert first slide to JPG)
def convert_pptx_first_slide_to_jpg(pptx_path):
    print(f"Converting {pptx_path} to JPG...")
    # Generate output path with same name but jpg extension
    base_name = os.path.splitext(os.path.basename(pptx_path))[0]
    output_dir = os.path.dirname(pptx_path)
    output_path = os.path.join(output_dir, base_name)
    
    # Use LibreOffice to convert
    subprocess.call([
        'soffice',
        '--headless',
        '--convert-to', 'jpg',
        '--outdir', output_dir,
        pptx_path
    ])
    
    # This will create JPGs for all slides, you'd need to keep only the first one
    return output_path + ".jpg"



# DESCRIPTION (generate description from overall content)

def extract_text_from_pptx(file_path):
    print(f"Extracting text from {file_path}...")
    if not os.path.exists(file_path):
        print(f"Error: File not found at {file_path}")
        return []
    
    try:
        prs = Presentation(file_path)
        slides_content = []
        
        for i, slide in enumerate(prs.slides):
            slide_text = []
            
            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            
            # Add slide content to results
            slides_content.append({
                "slide_number": i + 1,
                "content": "\n".join(slide_text)
            })
        
        return slides_content
    
    except Exception as e:
        print(f"Error extracting text from PowerPoint: {e}")
        return []



def generate_description_from_ppt(text_content):
    print(f"Generating description from text_content...")
    # result = ollama.generate("llama3.2", f"""
    # result = ollama.generate("gemma3", f"""
    result = ollama.generate("llama3.3", f"""

You are writing as Nicolas Deville, a senior sales director at Kaltura. Adopt his tone of voice and writing style, based on the following characteristics:

	•	Clear, structured, and concise: use bullet points or numbered lists when appropriate, and organise thoughts logically.
	•	Professional but personable: maintain a polite, human tone. Add light, personal touches when relevant (e.g., a friendly sign-off, reference to a location or recent event).
	•	Context-aware and analytical: reflect curiosity and insight, especially when analysing situations or communicating internally.
	•	Value-focused: highlight relevance and benefits for the reader, without being overly promotional.
	•	Assertive, but constructive: feel confident voicing concerns or driving action, while staying solutions-oriented.
	•	Multilingual-capable: write in English, German, or French depending on the recipient. Adapt tone and formality to suit the cultural context..

Write a very short description of the following PowerPoint presentation (text content extracted from the slides) as if I was writing it (ie describe the content to the person who will be watching it, in my voice):

---
{text_content}
---

Do not include any other text than the description.

""")
    
    output = '\n'.join(line.strip() for line in result.response.splitlines())

    return output


# TITLE (generate title from file name)
def generate_title_from_ppt_filename(file_path):
    print(f"Generating title from {file_path}...")
    
    # Extract the filename from the path
    filename = os.path.basename(file_path)
    
    # Remove the first 6 characters (YYMMDD timestamp) and strip
    title = filename[6:].strip()
    
    # Remove file extension if present
    title = os.path.splitext(title)[0]

    print(f"\t> Title: {title}")
    
    return title.strip('"\'')

# Extract title from first slide of PowerPoint
def generate_title_from_ppt_first_slide(file_path):
    print(f"Extracting title from first slide of {file_path}...")
    
    try:
        # Import pptx library for PowerPoint manipulation
        from pptx import Presentation
        
        # Load the presentation
        presentation = Presentation(file_path)
        
        # Get the first slide
        if len(presentation.slides) > 0:
            first_slide = presentation.slides[0]
            
            # Extract text from all shapes in the first slide
            title_text = ""
            for shape in first_slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    # Use the first non-empty text element as the title
                    title_text = shape.text.strip()
                    break
            
            if title_text:
                print(f"\t> Title from first slide: {title_text}")
                return title_text.strip('"\'')
            else:
                print(f"\t> No text found in first slide, falling back to filename")
                return generate_title_from_ppt_filename(file_path)
        else:
            print(f"\t> No slides found in presentation, falling back to filename")
            return generate_title_from_ppt_filename(file_path)
            
    except Exception as e:
        print(f"\t❌ Error extracting title from first slide: {str(e)}")
        print(f"\t> Falling back to filename-based title")
        return generate_title_from_ppt_filename(file_path)



# Upload to Kaltura
def upload_ppt_to_kaltura(file_path, channel_id=None, demo_mode=False, ppt_thumbnail_path=None, title_text=None, subtitle_text=None, MY_USER_SECRET=None, MY_ADMIN_SECRET=None, MY_PARTNER_ID=None): 

    # Kaltura service configuration
    config = KalturaConfiguration()
    config.serviceUrl = "https://www.kaltura.com/"
    client = KalturaClient(config)
    
    # Load the Document plugin
    client.loadPlugin('Document')
    
    expiry = 86400
    privileges = "disableentitlement"

    # Start session
    print(f"Starting session...")
    ks = client.session.start(MY_ADMIN_SECRET, MY_USER_SECRET, KalturaSessionType.ADMIN, MY_PARTNER_ID, expiry, privileges)
    client.setKs(ks)
    
    # Create upload token
    upload_token = KalturaUploadToken()
    upload_token = client.uploadToken.add(upload_token)
    
    # File to upload
    file_data = open(file_path, "rb")
    
    # Upload the file
    print(f"Uploading file...")
    resume = False
    final_chunk = True
    resume_at = -1
    result = client.uploadToken.upload(upload_token.id, file_data, resume, final_chunk, resume_at)

    # Create document entry instead of media entry
    document_entry = KalturaDocumentEntry()
    
    # Document Type - specify PowerPoint
    document_entry.documentType = KalturaDocumentType.DOCUMENT
    
    # Reference ID
    document_entry.referenceId = "auto-" + datetime.now().strftime('%Y%m%d%H%M%S')
    
    # ✅ TITLE
    if not demo_mode:
        document_entry.name = generate_title_from_ppt_filename(file_path)
    else:
        document_entry.name = title_text # Generated by AI and passed as argument
    # Remove any quotes at the beginning and end of title
    if document_entry.name:
        document_entry.name = document_entry.name.strip('"\'')

    # ✅ DESCRIPTION
    # Generate description from overall content using AI
    if not demo_mode:
        document_entry.description = generate_description_from_ppt(extract_text_from_pptx(file_path))
    else:
        document_entry.description = subtitle_text # Generated by AI and passed as argument
    
    # Set creator and owner information
    document_entry.userId = "nicolas.deville@kaltura.com"
    document_entry.creatorId = "nicolas.deville@kaltura.com"
    


    # ADD THE DOCUMENT ENTRY
    document_entry = client.baseEntry.add(document_entry)
    
    # Attach uploaded file to the document entry
    resource = KalturaUploadedFileTokenResource()
    resource.token = upload_token.id
    document_entry = client.baseEntry.addContent(document_entry.id, resource)
    
    # ✅ THUMBNAIL
    if not demo_mode:
        thumbnail_file_path = convert_pptx_first_slide_to_jpg(file_path)
    else:
        thumbnail_file_path = ppt_thumbnail_path # Generated by AI and passed as argument

    try:
        # Create an upload token for the thumbnail
        thumb_token = KalturaUploadToken()
        thumb_token = client.uploadToken.add(thumb_token)
        
        # Upload the thumbnail file
        thumb_file = open(thumbnail_file_path, 'rb')
        client.uploadToken.upload(thumb_token.id, thumb_file)
        thumb_file.close()
        
        # Create a new thumbAsset
        thumb_asset = KalturaThumbAsset()
        
        # Add the thumbnail to the document entry
        result_thumb = client.thumbAsset.add(document_entry.id, thumb_asset)
        
        # Create a resource from the upload token
        thumb_resource = KalturaUploadedFileTokenResource()
        thumb_resource.token = thumb_token.id
        
        # Set the content of the thumbnail
        client.thumbAsset.setContent(result_thumb.id, thumb_resource)
        
        # Set this thumbnail as the default for the document entry
        client.thumbAsset.setAsDefault(result_thumb.id)
        
        print(f"✅ Thumbnail successfully added & set as default")

    except Exception as e:
        print(f"❌ Error adding thumbnail: {e}")


    # ✅ CATEGORY (add to a Category so it appears in KMS Channel)
    category_entry = KalturaCategoryEntry()
    category_entry.entryId = document_entry.id
    category_entry.categoryId = channel_id
    
    try:
        client.categoryEntry.add(category_entry)
        print(f"✅ Successfully added entry to category ID: {channel_id}")
    except Exception as e:
        print(f"❌ Error adding to category: {e}")
        exit()
    
    document_url = f"https://nicolas.mediaspace.kaltura.com/media/{document_entry.id}"
    # Print the URL of the uploaded document
    print(f"\n🎉🎉🎉 Upload successful of {file_path} with title '{document_entry.name} at {document_url}\n\n")
    if not demo_mode:
        # Copy document URL to clipboard
        subprocess.run(['pbcopy'], input=document_url.encode('utf-8'))
        print("✅ Document URL copied to clipboard")

    file_data.close()

    return document_url
if __name__ == "__main__":

    KMS = "Pharma"

    # Credentials
    # Set up Kaltura credentials based on the selected KMS
    if KMS == "MY_KMS":
        USER_SECRET = os.getenv("user_secret")
        ADMIN_SECRET = os.getenv("admin_secret")
        PARTNER_ID = os.getenv("partner_id")
    elif KMS == "Pharma":
        USER_SECRET = os.getenv("pharma_user_secret")
        ADMIN_SECRET = os.getenv("pharma_admin_secret")
        PARTNER_ID = os.getenv("pharma_partner_id")
    elif KMS == "ABB":
        USER_SECRET = os.getenv("abb_user_secret")
        ADMIN_SECRET = os.getenv("abb_admin_secret")
        PARTNER_ID = os.getenv("abb_partner_id")
    else:
        print(f"❌ Error: Unknown KMS '{KMS}'. Please check your configuration.")
        exit()


    # file_path = "/Users/nic/Downloads/temp/test-ppt-upload.pptx"
    file_path = "/Users/nic/demo/pharma/pharma-demo.pptx"
    channel_id = 374884072

    print(f"\n\nUploading {file_path} to Kaltura with custom title, description, and a custom thumbnail...")
    upload_ppt_to_kaltura(file_path, channel_id=channel_id, demo_mode=True, MY_USER_SECRET=USER_SECRET, MY_ADMIN_SECRET=ADMIN_SECRET, MY_PARTNER_ID=PARTNER_ID)
    print(f"\n\n✅ File uploaded successfully to Kaltura")

    # End Chrono
    run_time = round((time.time() - start_time), 3)
    print(f'\n{os.path.basename(__file__)} finished in {round(run_time)}s at {datetime.now().strftime("%H:%M:%S")}.\n')
