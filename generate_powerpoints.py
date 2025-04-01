from datetime import datetime
import os
import sys
import time
import shutil
from pptx import Presentation
import requests
import ollama
import subprocess

"""
TODO:
- generate titles/subtitles with more randomness in the topic by passing a list of already generated titles/subtitles
"""

def generate_title_with_ollama(theme):

    result = ollama.generate("llama3.3", f"""

Generate a creative and concise title for a presentation part of the following theme: {theme}.
Only return the title, no other text.
This will be inserted into the title slide of a PowerPoint presentation.
                                                            
    """)

    output = '\n'.join(line.strip() for line in result.response.splitlines())

    print(f"generate_title_with_ollama > {output}\n\n")

    return output



def generate_subtitle_with_ollama(title):

    result = ollama.generate("llama3.3", f"""

Generate a creative and concise subtitle for a presentation with title: '{title}'.
Only return the subtitle, no other text.
This will be inserted into the subtitle text field in the title slide of a PowerPoint presentation.
                                                            
    """)

    output = '\n'.join(line.strip() for line in result.response.splitlines())

    print(f"generate_subtitle_with_ollama > {output}\n\n")

    return output



def generate_pptx(file_path, theme, num_presentations=1):
    """
    Duplicates a PowerPoint file and changes the text in the first slide using Ollama.
    
    Args:
        file_path (str): Path to the PowerPoint file to modify
        theme (str): Theme to use for generating new text
        num_presentations (int): Number of presentations to generate
    """
    # Check if file exists
    if not os.path.exists(file_path):
        print(f"Error: File '{file_path}' does not exist.")
        return
    
    # Check if file is a PowerPoint file
    if not file_path.lower().endswith('.pptx'):
        print(f"Error: File '{file_path}' is not a PowerPoint file (.pptx).")
        return
    
    generated_files = []
    
    for i in range(num_presentations):
        print(f"\n\n========= Creating presentation {i+1} of {num_presentations}\n")
        
        # Create a new file name with timestamp
        file_dir = os.path.dirname(file_path)
        file_name = os.path.basename(file_path)
        name_without_ext = os.path.splitext(file_name)[0]
        timestamp = datetime.now().strftime('%y%m%d%H%M%S')
        # Add index to filename to ensure uniqueness
        new_file_name = f"{name_without_ext}_{i+1:02d}.pptx"
        new_file_path = os.path.join(file_dir, new_file_name)
        
        # try:
            # First, make a copy of the original file
        shutil.copy2(file_path, new_file_path)
        print(f"Created copy: {new_file_path}")
        
        # Open the copied presentation
        prs = Presentation(new_file_path)
        
        # Check if there are slides
        if len(prs.slides) == 0:
            print("Warning: The presentation has no slides.")
            continue
        
        # Get the first slide
        slide = prs.slides[0]
        
        # Replace text in placeholders
        text_replaced = False
        generated_title = generate_title_with_ollama(theme)
        generated_subtitle = generate_subtitle_with_ollama(generated_title)
        
        # Track if we've found the first and second text shapes
        title_shape_found = False
        
        for shape in slide.shapes:
            if shape.has_text_frame:
                if not title_shape_found:
                    # First text shape gets the title
                    shape.text_frame.text = generated_title
                    title_shape_found = True
                    text_replaced = True
                else:
                    # Second text shape gets the subtitle
                    shape.text_frame.text = generated_subtitle
                    text_replaced = True
                
            # Handle tables if present
            elif hasattr(shape, "table"):
                for row in shape.table.rows:
                    for cell in row.cells:
                        cell.text = generate_subtitle_with_ollama(generated_title)
                        text_replaced = True

        if not text_replaced:
            print("Warning: No text frames found in the first slide.")
        
        # Save the modified presentation
        prs.save(new_file_path)
        
        generated_files.append(new_file_path)
        print(f"✅ Saved new presentation > {new_file_path}")
            
        # except Exception as e:
        #     print(f"Error modifying PowerPoint file: {e}")
    
    return generated_files


if __name__ == "__main__":

    # Start Chrono
    import time
    start_time = time.time()
    
    pptx_path = "/Users/nic/demo/pharma/pharma-demo.pptx"
    theme = "Pharmarceutical MSL & HCP Training presentations"
    num_presentations = 5
    
    generate_pptx(pptx_path, theme, num_presentations)
    
    
    run_time = time.time() - start_time
    if run_time < 1:
        print(f"\nFinished creating {num_presentations} presentations in {round(run_time*1000)}ms at {datetime.now().strftime('%H:%M:%S')}.\n")
    elif run_time < 60:
        print(f"\nFinished creating {num_presentations} presentations in {round(run_time)}s at {datetime.now().strftime('%H:%M:%S')}.\n")
    else:
        print(f"\nFinished creating {num_presentations} presentations in {round(run_time/60)}mins at {datetime.now().strftime('%H:%M:%S')}.\n")
